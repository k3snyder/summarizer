#!/usr/bin/env python
"""
Creates test_presentation.pptx with 6 slides for testing PPTX extraction.

Slide 1: Title slide with 'Test Presentation' title (no notes)
Slide 2: Bullet points + speaker notes ('Remember to explain this')
Slide 3: Table (3 rows x 2 columns: Header1/Header2, A/B, C/D)
Slide 4: Image-only slide (embed a small test image, no text shapes)
Slide 5: Empty/blank slide
Slide 6: Two tables side by side
"""

import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def create_test_image() -> bytes:
    """Create a small test image (100x100 red square) as bytes."""
    img = Image.new("RGB", (100, 100), color=(255, 0, 0))  # type: ignore[arg-type]
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer.read()


def create_test_presentation(output_path: Path) -> None:
    """Create test presentation with 6 slides."""
    prs = Presentation()

    # Slide 1: Title slide with 'Test Presentation' title (no notes)
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide1 = prs.slides.add_slide(slide_layout)
    title = slide1.shapes.title
    if title is not None:
        title.text = "Test Presentation"
    subtitle = slide1.placeholders[1]
    if subtitle is not None and hasattr(subtitle, "text"):
        subtitle.text = "Subtitle text for testing"  # type: ignore[union-attr]

    # Slide 2: Bullet points + speaker notes
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide2 = prs.slides.add_slide(slide_layout)
    title = slide2.shapes.title
    if title is not None:
        title.text = "Slide with Bullet Points"
    body_shape = slide2.placeholders[1]
    if body_shape is None or not hasattr(body_shape, "text_frame"):
        raise RuntimeError("Expected content placeholder in slide layout")
    tf = body_shape.text_frame  # type: ignore[union-attr]
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Third bullet point"
    p.level = 0

    # Add speaker notes
    notes_slide = slide2.notes_slide
    notes_tf = notes_slide.notes_text_frame
    if notes_tf is not None:
        notes_tf.text = "Remember to explain this"

    # Slide 3: Table (3 rows x 2 columns)
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide3 = prs.slides.add_slide(slide_layout)

    # Add title text box
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Slide with Table"

    # Add table
    rows, cols = 3, 2
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(2)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    # Set table data
    table.cell(0, 0).text = "Header1"
    table.cell(0, 1).text = "Header2"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    table.cell(2, 0).text = "C"
    table.cell(2, 1).text = "D"

    # Slide 4: Image-only slide (no text shapes except the image)
    slide_layout = prs.slide_layouts[6]  # Blank layout (completely blank)
    slide4 = prs.slides.add_slide(slide_layout)

    # Create and add test image
    img_bytes = create_test_image()
    img_stream = io.BytesIO(img_bytes)
    left = Inches(3)
    top = Inches(2)
    width = Inches(4)
    slide4.shapes.add_picture(img_stream, left, top, width=width)

    # Slide 5: Empty/blank slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    _ = prs.slides.add_slide(slide_layout)
    # No shapes added - completely empty

    # Slide 6: Two tables side by side
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide6 = prs.slides.add_slide(slide_layout)

    # Add title text box
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide6.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Slide with Two Tables"

    # Table 1 (left)
    rows, cols = 2, 2
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(4)
    height = Inches(1.5)
    table1 = slide6.shapes.add_table(rows, cols, left, top, width, height).table
    table1.cell(0, 0).text = "T1-R1C1"
    table1.cell(0, 1).text = "T1-R1C2"
    table1.cell(1, 0).text = "T1-R2C1"
    table1.cell(1, 1).text = "T1-R2C2"

    # Table 2 (right)
    left = Inches(5.5)
    table2 = slide6.shapes.add_table(rows, cols, left, top, width, height).table
    table2.cell(0, 0).text = "T2-R1C1"
    table2.cell(0, 1).text = "T2-R1C2"
    table2.cell(1, 0).text = "T2-R2C1"
    table2.cell(1, 1).text = "T2-R2C2"

    # Save presentation
    prs.save(str(output_path))
    print(f"Created test presentation: {output_path}")
    print(f"  - Slide 1: Title slide")
    print(f"  - Slide 2: Bullet points + speaker notes")
    print(f"  - Slide 3: Table (3x2)")
    print(f"  - Slide 4: Image only")
    print(f"  - Slide 5: Empty/blank")
    print(f"  - Slide 6: Two tables side by side")


if __name__ == "__main__":
    output_path = Path(__file__).parent / "test_presentation.pptx"
    create_test_presentation(output_path)
