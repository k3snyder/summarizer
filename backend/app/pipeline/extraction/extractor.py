"""
DocumentExtractor - Async document extraction service.

Handles PDF and text file extraction with progress callbacks.
Uses asyncio.to_thread() for blocking operations.
"""

import asyncio
import base64
import logging
import re
from io import BytesIO
from pathlib import Path
from typing import Awaitable, Callable, Optional

from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.pipeline.extraction.schemas import (
    ExtractedDocument,
    ExtractedPage,
    ExtractionConfig,
)

logger = logging.getLogger(__name__)

# Type alias for progress callback
ProgressCallback = Callable[[int, int, str], Awaitable[None]]


class DocumentExtractor:
    """Async document extraction service for PDF, text, and PPTX files."""

    SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".pptx"}

    def __init__(self):
        """Initialize the document extractor."""
        pass

    async def extract(
        self,
        file_path: Path,
        config: ExtractionConfig,
        progress_callback: Optional[ProgressCallback] = None,
    ) -> ExtractedDocument:
        """
        Extract content from a document file.

        Args:
            file_path: Path to the document (PDF, TXT, or MD)
            config: Extraction configuration
            progress_callback: Optional async callback for progress updates
                Signature: (current: int, total: int, message: str) -> Awaitable[None]

        Returns:
            ExtractedDocument with all extracted content

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file type is unsupported
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = file_path.suffix.lower()
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {extension}. "
                f"Supported: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

        if extension == ".pdf":
            return await self._extract_pdf(file_path, config, progress_callback)
        elif extension == ".pptx":
            return await self._extract_pptx(file_path, config, progress_callback)
        else:
            return await self._extract_text(file_path, config, progress_callback)

    async def _extract_text(
        self,
        file_path: Path,
        config: ExtractionConfig,
        progress_callback: Optional[ProgressCallback] = None,
    ) -> ExtractedDocument:
        """Extract content from a text/markdown file."""
        # Read file content in thread to avoid blocking
        content = await asyncio.to_thread(self._read_text_file, file_path)

        # Split into chunks using langchain splitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.chunk_size,
            chunk_overlap=config.chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""],
        )
        chunks = splitter.split_text(content)
        total_chunks = len(chunks)

        if progress_callback:
            await progress_callback(0, total_chunks, f"Starting text extraction: {total_chunks} chunks")

        # Build pages from chunks
        pages: list[ExtractedPage] = []
        document_id = f"doc_{file_path.stem}"

        for idx, chunk_text in enumerate(chunks, start=1):
            page = ExtractedPage(
                page_number=idx,
                chunk_id=f"chunk_{idx}",
                doc_title=file_path.name,
                text=chunk_text,
                tables=[],
                image_base64=None,
                image_text=None,
                image_classifier=None,
            )
            pages.append(page)

            if progress_callback:
                await progress_callback(idx, total_chunks, f"Processed chunk {idx} of {total_chunks}")

        return ExtractedDocument(
            document_id=document_id,
            filename=file_path.name,
            total_pages=total_chunks,
            pages=pages,
            metadata={
                "source_type": "text",
                "file_type": file_path.suffix,
                "chunk_size": config.chunk_size,
                "chunk_overlap": config.chunk_overlap,
            },
        )

    def _read_text_file(self, file_path: Path) -> str:
        """Read text file content (blocking, called via to_thread)."""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    async def _extract_pdf(
        self,
        file_path: Path,
        config: ExtractionConfig,
        progress_callback: Optional[ProgressCallback] = None,
    ) -> ExtractedDocument:
        """Extract content from a PDF file."""
        # Open PDF and get page count (blocking operation)
        pdf_info = await asyncio.to_thread(self._get_pdf_info, file_path)
        total_pages = pdf_info["total_pages"]

        if progress_callback:
            await progress_callback(0, total_pages, f"Starting PDF extraction: {total_pages} pages")

        pages: list[ExtractedPage] = []
        document_id = f"doc_{file_path.stem}"

        # Process each page
        for page_num in range(1, total_pages + 1):
            page_data = await asyncio.to_thread(
                self._extract_pdf_page,
                file_path,
                page_num,
                config,
            )

            page = ExtractedPage(
                page_number=page_num,
                chunk_id=f"chunk_{page_num}",
                doc_title=file_path.name,
                text=page_data["text"],
                tables=page_data["tables"],
                image_base64=page_data["image_base64"],
                image_text=None,
                image_classifier=None,
            )
            pages.append(page)

            if progress_callback:
                await progress_callback(page_num, total_pages, f"Processed page {page_num} of {total_pages}")

        return ExtractedDocument(
            document_id=document_id,
            filename=file_path.name,
            total_pages=total_pages,
            pages=pages,
            metadata={
                "source_type": "pdf",
                "file_type": ".pdf",
            },
        )

    def _get_pdf_info(self, file_path: Path) -> dict:
        """Get PDF info (blocking, called via to_thread)."""
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            return {"total_pages": len(pdf.pages)}

    def _extract_pdf_page(
        self,
        file_path: Path,
        page_number: int,
        config: ExtractionConfig,
    ) -> dict:
        """
        Extract content from a single PDF page (blocking, called via to_thread).

        Returns dict with text, tables, and image_base64.
        """
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            page = pdf.pages[page_number - 1]

            # Extract text
            text = self._extract_text_from_page(page)

            # Extract tables
            tables: list[list[list[str]]] = []
            if not config.skip_tables and not config.text_only:
                tables = self._extract_tables_from_page(page)

            # Extract page image
            image_base64: Optional[str] = None
            if not config.skip_images and not config.text_only:
                image_base64 = self._extract_image_from_page(
                    file_path, page_number, config.pdf_image_dpi
                )

            return {
                "text": text,
                "tables": tables,
                "image_base64": image_base64,
            }

    def _extract_text_from_page(self, page) -> str:
        """
        Extract text from a pdfplumber page.
        Falls back to OCR if no text is found.
        """
        try:
            text = page.extract_text()
            if text and text.strip():
                return text

            # Fallback to OCR
            try:
                import pytesseract

                page_img = page.to_image(resolution=150)
                pil_img = page_img.original
                ocr_text = pytesseract.image_to_string(pil_img)
                return ocr_text.strip() if ocr_text else ""
            except ImportError:
                logger.warning("pytesseract not available for OCR fallback")
                return ""
            except Exception as e:
                logger.error(f"OCR failed: {e}")
                return ""
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return ""

    def _extract_tables_from_page(self, page) -> list[list[list[str]]]:
        """
        Extract tables from a pdfplumber page using multiple strategies.
        Returns list of tables, where each table is a list of rows (list of cells).
        """
        tables: list[list[list[str]]] = []
        try:
            strategies = [
                {
                    "name": "lines_strict",
                    "settings": {
                        "horizontal_strategy": "lines_strict",
                        "vertical_strategy": "lines_strict",
                        "intersection_y_tolerance": 3,
                    },
                },
                {
                    "name": "lines",
                    "settings": {
                        "horizontal_strategy": "lines",
                        "vertical_strategy": "lines",
                        "snap_tolerance": 3,
                        "join_tolerance": 3,
                        "edge_min_length": 3,
                        "min_words_vertical": 3,
                        "min_words_horizontal": 1,
                    },
                },
                {
                    "name": "text",
                    "settings": {
                        "horizontal_strategy": "text",
                        "vertical_strategy": "text",
                        "snap_tolerance": 3,
                        "intersection_y_tolerance": 3,
                    },
                },
            ]

            found_valid_tables = False
            for strategy in strategies:
                if found_valid_tables:
                    break

                found_tables = page.find_tables(table_settings=strategy["settings"])
                if not found_tables:
                    continue

                for table in found_tables:
                    try:
                        extracted = table.extract()
                        if not self._is_valid_table(extracted):
                            continue

                        # Clean and convert to list of lists
                        cleaned_table = []
                        for row in extracted:
                            cleaned_row = []
                            for cell in row:
                                if cell is None or str(cell).strip() == "":
                                    cleaned_row.append("")
                                else:
                                    cleaned_row.append(str(cell).strip())
                            cleaned_table.append(cleaned_row)

                        if len(cleaned_table) >= 2:
                            tables.append(cleaned_table)
                            found_valid_tables = True

                    except Exception as e:
                        logger.debug(f"Error extracting table: {e}")
                        continue

            # Try chart data extraction if no regular tables found
            if not tables:
                chart_tables = self._extract_chart_data(page)
                tables.extend(chart_tables)

        except Exception as e:
            logger.error(f"Table extraction failed: {e}")

        return tables

    def _is_valid_table(self, table_data) -> bool:
        """Check if extracted table data represents a real table."""
        if not table_data or len(table_data) < 2:
            return False

        total_cells = 0
        non_empty_cells = 0
        for row in table_data:
            for cell in row:
                total_cells += 1
                if cell and str(cell).strip() and cell != "None":
                    non_empty_cells += 1

        if total_cells == 0 or (non_empty_cells / total_cells) < 0.3:
            return False

        cols_with_data = 0
        for col_idx in range(len(table_data[0])):
            col_has_data = False
            for row in table_data:
                if col_idx < len(row) and row[col_idx] and str(row[col_idx]).strip():
                    col_has_data = True
                    break
            if col_has_data:
                cols_with_data += 1

        return cols_with_data >= 2

    def _extract_chart_data(self, page) -> list[list[list[str]]]:
        """Extract data from charts/figures that contain percentage data."""
        tables: list[list[list[str]]] = []
        try:
            text = page.extract_text()
            if not text:
                return tables

            percentage_pattern = r"([A-Za-z][A-Za-z\s\-,/]+?)\s+(\d+%)"
            matches = re.findall(percentage_pattern, text)

            data_items = []
            for item, percentage in matches:
                item = item.strip()
                if len(item) > 5 and not any(
                    skip in item.lower() for skip in ["share of", "surveyed", "source"]
                ):
                    data_items.append([item, percentage])

            if len(data_items) >= 3:
                seen = set()
                unique_items = []
                for item in data_items:
                    key = (item[0], item[1])
                    if key not in seen:
                        seen.add(key)
                        unique_items.append(item)

                # Create table with header row
                table = [["Item", "Percentage"]] + unique_items
                tables.append(table)

        except Exception as e:
            logger.error(f"Chart data extraction failed: {e}")

        return tables

    def _extract_image_from_page(
        self,
        file_path: Path,
        page_number: int,
        dpi: int,
    ) -> Optional[str]:
        """
        Extract full page as base64-encoded JPEG image.
        """
        try:
            from pdf2image import convert_from_path

            images = convert_from_path(
                str(file_path),
                first_page=page_number,
                last_page=page_number,
                dpi=dpi,
            )

            if images:
                return self._pil_image_to_base64(images[0])
            return None

        except ImportError:
            logger.warning("pdf2image not available for image extraction")
            return None
        except Exception as e:
            logger.error(f"Image extraction failed for page {page_number}: {e}")
            return None

    def _pil_image_to_base64(self, img) -> str:
        """Convert PIL Image to base64-encoded JPEG string."""
        buffered = BytesIO()
        img.save(buffered, format="JPEG")
        return base64.b64encode(buffered.getvalue()).decode("utf-8")

    async def _extract_pptx(
        self,
        file_path: Path,
        config: ExtractionConfig,
        progress_callback: Optional[ProgressCallback] = None,
    ) -> ExtractedDocument:
        """Extract content from a PowerPoint file."""
        from pptx import Presentation

        prs = await asyncio.to_thread(Presentation, str(file_path))
        total_slides = len(prs.slides)

        if progress_callback:
            await progress_callback(0, total_slides, "Starting PPTX extraction...")

        # Render all slide images first (requires LibreOffice)
        slide_images: dict[int, str] = {}
        if not config.skip_images and not config.text_only:
            slide_images = await asyncio.to_thread(
                self._render_pptx_slides, file_path, config
            )

        pages: list[ExtractedPage] = []
        # For PPTX, skip_tables config controls speaker notes (UI shows "Skip Slide Notes")
        skip_notes = config.skip_tables or config.text_only
        # skip_pptx_tables controls table extraction (UI shows "Skip Slide Tables")
        skip_slide_tables = config.skip_pptx_tables or config.text_only
        for slide_num, slide in enumerate(prs.slides, start=1):
            text, tables = self._extract_pptx_slide(slide, skip_notes, skip_slide_tables)

            page = ExtractedPage(
                page_number=slide_num,
                chunk_id=f"chunk_{slide_num}",
                doc_title=file_path.name,
                text=text,
                tables=tables,
                image_base64=slide_images.get(slide_num),
                image_text=None,
                image_classifier=None,
            )
            pages.append(page)

            if progress_callback:
                await progress_callback(
                    slide_num, total_slides, f"Slide {slide_num}/{total_slides}"
                )

        return ExtractedDocument(
            document_id=f"doc_{file_path.stem}",
            filename=file_path.name,
            total_pages=total_slides,
            pages=pages,
            metadata={"source_type": "pptx", "file_type": ".pptx"},
        )

    def _extract_pptx_slide(
        self, slide, skip_notes: bool, skip_tables: bool
    ) -> tuple[str, list[list[list[str]]]]:
        """
        Extract text, notes, and tables from a single slide.

        Args:
            slide: pptx Slide object
            skip_notes: If True, skip speaker notes extraction
            skip_tables: If True, skip table extraction

        Returns:
            Tuple of (text, tables) where tables is list of 2D arrays
        """
        # Extract ALL text from slide XML (catches SmartArt, groups, everything)
        text_parts = self._extract_all_slide_text(slide)

        # Extract tables (unless skip_tables is True)
        tables: list[list[list[str]]] = []
        if not skip_tables:
            tables = self._extract_tables_from_slide(slide)

        # Append speaker notes with marker (unless skip_notes is True)
        if not skip_notes and slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                text_parts.append("")
                text_parts.append("[Speaker Notes]")
                text_parts.append(notes_frame.text.strip())

        return "\n".join(text_parts), tables

    def _extract_all_slide_text(self, slide) -> list[str]:
        """
        Extract text from a slide by parsing XML directly, EXCLUDING table content.

        This captures text from all shape types including SmartArt, grouped shapes,
        and any other elements that python-pptx doesn't expose via has_text_frame.
        Table text is excluded to avoid duplication (tables are extracted separately).

        Args:
            slide: pptx Slide object

        Returns:
            List of text strings found in the slide (excluding tables)
        """
        text_parts: list[str] = []
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        tbl_tag = f"{ns}tbl"
        t_tag = f"{ns}t"

        slide_xml = slide._element

        # Find all text elements, but skip those inside tables
        for text_elem in slide_xml.iter(t_tag):
            # Check if this text element is inside a table by walking up the tree
            if self._is_inside_table(text_elem, tbl_tag):
                continue
            if text_elem.text and text_elem.text.strip():
                text_parts.append(text_elem.text.strip())

        return text_parts

    def _is_inside_table(self, element, tbl_tag: str) -> bool:
        """Check if an XML element is nested inside a table element."""
        parent = element.getparent()
        while parent is not None:
            if parent.tag == tbl_tag:
                return True
            parent = parent.getparent()
        return False

    def _extract_tables_from_slide(self, slide) -> list[list[list[str]]]:
        """
        Extract tables from a slide by parsing XML directly.

        This approach captures ALL tables including those in SmartArt or grouped
        shapes, and properly handles merged cells by only extracting content from
        merge-origin cells (skipping continuation cells marked with hMerge/vMerge).

        Args:
            slide: pptx Slide object

        Returns:
            List of tables, where each table is a 2D array of strings
        """
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        tables: list[list[list[str]]] = []

        slide_xml = slide._element

        # Find all table elements in the slide XML
        for tbl_elem in slide_xml.iter(f"{ns}tbl"):
            table_data = self._extract_table_from_xml(tbl_elem, ns)
            if table_data:
                tables.append(table_data)

        return tables

    def _extract_table_from_xml(
        self, tbl_elem, ns: str
    ) -> list[list[str]]:
        """
        Extract table data from an XML table element.

        Handles merged cells by:
        - Extracting content only from merge-origin cells (gridSpan/rowSpan)
        - Skipping continuation cells (hMerge="1" or vMerge="1")

        Args:
            tbl_elem: lxml Element for <a:tbl>
            ns: XML namespace string

        Returns:
            2D array of cell strings (rows x columns)
        """
        rows: list[list[str]] = []

        # Find all table rows
        for tr_elem in tbl_elem.findall(f"{ns}tr"):
            row_cells: list[str] = []

            # Find all cells in this row
            for tc_elem in tr_elem.findall(f"{ns}tc"):
                # Check for merge continuation flags (skip these cells)
                tc_pr = tc_elem.find(f"{ns}tcPr")
                if tc_pr is not None:
                    # hMerge="1" means this cell is a horizontal merge continuation
                    # vMerge="1" means this cell is a vertical merge continuation
                    h_merge = tc_pr.get("hMerge")
                    v_merge = tc_pr.get("vMerge")
                    if h_merge == "1" or v_merge == "1":
                        # Skip continuation cells (content is in the origin cell)
                        continue

                # Extract all text from this cell
                cell_text = self._extract_cell_text(tc_elem, ns)
                row_cells.append(cell_text)

            if row_cells:
                rows.append(row_cells)

        return rows

    def _extract_cell_text(self, tc_elem, ns: str) -> str:
        """
        Extract all text from a table cell element.

        Args:
            tc_elem: lxml Element for <a:tc>
            ns: XML namespace string

        Returns:
            Concatenated text content of the cell
        """
        text_parts: list[str] = []

        # Find all <a:t> text elements within the cell
        for t_elem in tc_elem.iter(f"{ns}t"):
            if t_elem.text:
                text_parts.append(t_elem.text)

        return " ".join(text_parts).strip()

    def _render_pptx_slides(
        self, file_path: Path, config: ExtractionConfig
    ) -> dict[int, str]:
        """
        Render slides to JPEG images using LibreOffice headless + pdf2image.

        Args:
            file_path: Path to the PPTX file
            config: Extraction config (for DPI setting)

        Returns:
            Dict mapping slide number (1-based) to base64 JPEG string
        """
        import subprocess
        import tempfile

        from pdf2image import convert_from_path

        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                # Step 1: Convert PPTX to PDF via LibreOffice
                result = subprocess.run(
                    [
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        tmpdir,
                        str(file_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                )

                if result.returncode != 0:
                    logger.warning(f"LibreOffice conversion failed: {result.stderr}")
                    return {}

                # Step 2: Convert PDF pages to images
                pdf_path = Path(tmpdir) / f"{file_path.stem}.pdf"
                if not pdf_path.exists():
                    logger.warning(f"PDF not generated at {pdf_path}")
                    return {}

                images = convert_from_path(str(pdf_path), dpi=config.pdf_image_dpi)

                # Step 3: Encode each slide as base64 JPEG
                slide_images: dict[int, str] = {}
                for i, img in enumerate(images, start=1):
                    slide_images[i] = self._pil_image_to_base64(img)

                return slide_images

        except FileNotFoundError:
            logger.warning(
                "LibreOffice not found. Install it:\n"
                "  macOS: brew install --cask libreoffice\n"
                "  Linux: sudo apt-get install libreoffice"
            )
            return {}
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice conversion timed out (>120s)")
            return {}
